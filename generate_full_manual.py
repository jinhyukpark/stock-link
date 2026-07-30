from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def create_slide(prs, title, path, purpose, content, usage, icon_text, is_section=False):
    if is_section:
        slide_layout = prs.slide_layouts[0] # Title layout
        slide = prs.slides.add_slide(slide_layout)
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(30, 58, 138) # blue-900 (Section divider color)
        
        title_shape = slide.shapes.title
        title_shape.text = f"{icon_text} {title}"
        title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        title_shape.text_frame.paragraphs[0].font.size = Pt(48)
        title_shape.text_frame.paragraphs[0].font.bold = True
        
        subtitle = slide.placeholders[1]
        subtitle.text = purpose
        subtitle.text_frame.paragraphs[0].font.color.rgb = RGBColor(191, 219, 254) # blue-200
        return

    slide_layout = prs.slide_layouts[5] # Blank layout with title
    slide = prs.slides.add_slide(slide_layout)
    
    # Background color (dark theme)
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(15, 23, 42) # slate-900

    # Header Bar
    shapes = slide.shapes
    header_bar = shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1)
    )
    header_bar.fill.solid()
    header_bar.fill.fore_color.rgb = RGBColor(30, 41, 59) # slate-800
    header_bar.line.color.rgb = RGBColor(30, 41, 59)

    # Title
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.6))
    tf = txBox.text_frame
    tf.text = f"{icon_text} {title}"
    tf.paragraphs[0].font.size = Pt(24)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255) # white

    # Content Boxes
    left_margin = Inches(0.5)
    width = Inches(9)
    
    # 1. 화면 경로
    txBox = slide.shapes.add_textbox(left_margin, Inches(1.2), width, Inches(0.5))
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = "📍 화면 경로: "
    p.font.bold = True
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(96, 165, 250) # blue-400
    p.text += path
    p.font.bold = False
    p.font.color.rgb = RGBColor(203, 213, 225) # slate-300

    # 2. 화면 목적
    txBox = slide.shapes.add_textbox(left_margin, Inches(1.6), width, Inches(0.8))
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = "🎯 화면 목적 / 설명"
    p.font.bold = True
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(255, 255, 255)
    
    p2 = tf.add_paragraph()
    p2.text = purpose
    p2.font.size = Pt(12)
    p2.font.color.rgb = RGBColor(203, 213, 225)
    p2.space_before = Pt(4)

    # 3. 표출 내용
    txBox = slide.shapes.add_textbox(left_margin, Inches(2.6), Inches(4.3), Inches(4))
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = "📋 화면에 보이는 내용"
    p.font.bold = True
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(255, 255, 255)
    
    for item in content:
        p_item = tf.add_paragraph()
        p_item.text = f"• {item}"
        p_item.font.size = Pt(11)
        p_item.font.color.rgb = RGBColor(203, 213, 225)
        p_item.space_before = Pt(4)

    # 4. 사용 방법
    txBox = slide.shapes.add_textbox(Inches(5.2), Inches(2.6), Inches(4.3), Inches(4))
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = "💡 이렇게 사용하세요"
    p.font.bold = True
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(255, 255, 255)
    
    for i, item in enumerate(usage):
        p_item = tf.add_paragraph()
        p_item.text = f"{(i+1)}️⃣ {item}"
        p_item.font.size = Pt(11)
        p_item.font.color.rgb = RGBColor(203, 213, 225)
        p_item.space_before = Pt(4)

def create_manual():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)

    # 1. Title Slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(15, 23, 42)
    
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "StockLink 전체 화면 사용자 매뉴얼"
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    title.text_frame.paragraphs[0].font.size = Pt(40)
    title.text_frame.paragraphs[0].font.bold = True
    
    subtitle.text = "모든 탭, 팝업, 상세 화면을 포함한 완전판 안내서"
    subtitle.text_frame.paragraphs[0].font.color.rgb = RGBColor(148, 163, 184)
    
    # 2. 목차
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)
    background = slide.background
    background.fill.solid()
    background.fill.fore_color.rgb = RGBColor(15, 23, 42)
    
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(4.5))
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = "📑 스캔된 전체 화면 목차"
    p.font.bold = True
    p.font.size = Pt(24)
    p.font.color.rgb = RGBColor(255, 255, 255)
    
    toc = [
        "1. 홈(대시보드)",
        "2. 시장현황",
        "3. 인사이트 (마켓맵/뉴스/소셜/테마/비교 탭 및 모달)",
        "4. 종목 분석 (12개 개별 상세 탭)",
        "5. 관계망(온톨로지)",
        "6. 상승탄력 분석",
        "7. 커뮤니티 (목록, 모달)",
        "8. 마이페이지 / 구독 / 기타"
    ]
    
    for item in toc:
        p_item = tf.add_paragraph()
        p_item.text = item
        p_item.font.size = Pt(14)
        p_item.font.color.rgb = RGBColor(203, 213, 225)
        p_item.space_before = Pt(6)

    # --- 1. 홈 ---
    create_slide(prs, "1. 홈 (대시보드)", "홈 /", "섹션 커버", [], [], "🏠", is_section=True)
    
    create_slide(prs, "홈 메인 화면", "홈 (/) > 대시보드",
        "StockLink 접속 시 가장 먼저 만나는 화면으로, 현재 시장 요약과 내 관심 종목 상태를 한눈에 파악합니다.",
        ["주요 지수 요약 (KOSPI, 코스닥)", "공포/탐욕 지수 게이지 (시장의 심리 상태)", "글로벌 증시 및 시장 히트맵 (분야별 상승/하락 비율)", "오늘의 상승/하락 특징주", "트렌드 워드 클라우드 (오늘 가장 많이 언급된 단어들)"],
        ["공포/탐욕 지수를 보고 현재 시장 참가자들이 겁을 내는지, 욕심을 내는지 파악합니다.", "히트맵에서 빨간색(상승)과 파란색(하락) 비중을 보며 어떤 분야가 강세인지 살핍니다.", "워드 클라우드의 큰 글씨를 통해 오늘 시장의 핵심 키워드를 확인합니다."], "📊")

    # --- 2. 시장 현황 ---
    create_slide(prs, "2. 시장 현황", "시장 현황 /market", "섹션 커버", [], [], "📈", is_section=True)
    
    create_slide(prs, "시장 현황 대시보드", "메뉴 > 시장 현황 (/market)",
        "국내외 시장 전체가 현재 어떤 방향으로 흘러가고 있는지 깊이 있게 분석해주는 화면입니다.",
        ["시장 신호등: 시장의 과열(너무 비쌈) 또는 침체(너무 쌈) 상태를 알려주는 지표", "투자자 동향: 기관, 외국인, 개인이 주식을 사고 파는 흐름 차트", "주요 뉴스 요약 목록"],
        ["상단의 '시장 신호등'을 보고 지금 매수를 할지 보수적으로 지켜볼지 감을 잡습니다.", "외국인과 기관 그래프가 위로 향하는지 아래로 향하는지 확인하여 수급을 파악합니다."], "🚦")

    # --- 3. 인사이트 ---
    create_slide(prs, "3. 인사이트 (테마/소셜 분석)", "인사이트 /insight", "섹션 커버", [], [], "💡", is_section=True)
    
    create_slide(prs, "인사이트 > 마켓 맵 탭", "메뉴 > 인사이트 > 마켓 맵 탭",
        "전체 시장을 네모난 블록 형태로 시각화하여 어느 섹터가 오르고 내리는지 직관적으로 보여줍니다.",
        ["섹터별/종목별 등락률 블록(트리맵)", "상승(빨강), 하락(파랑) 색상 대비", "선택 섹터 상세 정보 패널"],
        ["화면의 블록 색상과 크기를 통해 가장 많이 오르거나 비중이 큰 섹터를 한눈에 파악합니다.", "특정 블록에 마우스를 올려 세부 종목들의 등락률을 확인합니다."], "🗺️")

    create_slide(prs, "인사이트 > 소셜 분석 탭", "메뉴 > 인사이트 > 소셜 분석 탭",
        "유명 인사나 SNS에서 주식에 대해 어떤 이야기들이 오가는지 요약하여 시장 영향도를 보여줍니다.",
        ["날짜 선택기 (과거 리포트 조회 가능)", "주요 하이라이트 요약 텍스트", "인플루언서 발언 및 긍정/부정(별점) 평가 테이블", "상승/하락 기대 종목 요약 카드", "긍정/부정 섹터별 파이 차트"],
        ["우측 상단에서 날짜를 선택합니다.", "테이블의 '발언 요약'을 통해 오늘 전문가들이 무슨 말을 했는지 훑어봅니다.", "파이 차트를 보며 긍정적인 이야기가 많은 섹터를 찾아봅니다."], "🗣️")

    create_slide(prs, "[팝업] 소셜 분석 > 발언 상세 및 관련 종목", "소셜 분석 탭 > 테이블 내 특정 인플루언서 행 클릭",
        "특정 인플루언서의 발언 원문과 그로 인해 영향을 받는 구체적인 긍정/부정 종목을 상세히 확인하는 팝업입니다.",
        ["발언 원문 및 플랫폼 정보", "📈 긍정 영향 종목 리스트 (해당 종목이 상승탄력을 받을 사유 포함)", "📉 부정 영향 종목 리스트 (하락 우려 사유 포함)", "국내 주식(이동 화살표 표출), 해외 주식('해외' 배지 표출)", "원문 보기 버튼"],
        ["테이블에서 관심 있는 인물(예: 일론 머스크)을 클릭합니다.", "팝업이 뜨면 발언 원문을 읽어 맥락을 파악합니다.", "아래의 긍정/부정 종목 리스트를 확인하며, 국내 종목의 경우 클릭하여 종목 상세 화면으로 이동할 수 있습니다."], "🔍")

    create_slide(prs, "인사이트 > 테마 분석 탭", "메뉴 > 인사이트 > 테마 분석 탭",
        "현재 주식 시장에서 가장 핫한 테마(예: AI, 2차전지 등)와 관련 종목들을 그룹지어 분석합니다.",
        ["오늘의 주도 테마 리스트 및 등락률", "테마별 주도주(가장 많이 오른 주식) 정보", "테마 트렌드 차트"],
        ["가장 등락률이 높은 상위 테마를 확인하여 돈이 어디로 몰리고 있는지 파악합니다.", "관심 있는 테마 카드를 클릭하여 상세 팝업을 엽니다."], "🗂️")

    create_slide(prs, "[팝업] 테마 상세 분석", "테마 분석 탭 > 특정 테마 카드 클릭",
        "선택한 테마의 상세 설명과 그 테마에 속한 모든 관련 종목을 표 형태로 확인하는 모달입니다.",
        ["테마 개요 및 오늘 상승한 이유", "테마 내 편입 종목 전체 리스트 (종목명, 현재가, 등락률)", "차트 및 관련 뉴스 요약"],
        ["팝업이 뜨면 해당 테마가 오늘 왜 오르는지 요약 설명을 읽습니다.", "아래 종목 리스트에서 가장 많이 오른 대장주를 찾아봅니다.", "원하는 종목을 클릭해 종목 분석 상세로 넘어갑니다."], "📑")

    create_slide(prs, "인사이트 > 주요 종목 비교 탭", "메뉴 > 인사이트 > 주요 종목 비교 탭",
        "관심 있는 여러 주식들을 한 화면에 올려놓고 실적, 투자지표 등을 나란히 비교합니다.",
        ["종목 추가 검색창", "비교 테이블 (수익률, 시가총액, PER, 영업이익률 등 가로 나열)", "항목별 하이라이트 (가장 수치가 좋은 곳에 색상 표시)"],
        ["'종목 추가' 버튼을 눌러 비교하고 싶은 주식(예: 삼성전자, SK하이닉스)을 차례로 넣습니다.", "테이블을 가로로 보면서 어떤 회사가 수익률이나 실적이 더 좋은지 비교합니다."], "⚖️")

    # --- 4. 종목 분석 (상세 탭들) ---
    create_slide(prs, "4. 종목 상세 분석 (12개 탭)", "종목 상세 /stock", "섹션 커버", [], [], "🏢", is_section=True)

    create_slide(prs, "종목 분석 > 상단 공통 영역", "메뉴 > 종목 분석 (/stock)",
        "검색한 특정 기업의 가장 핵심적인 현재 가격 상태를 보여주는 고정 영역입니다.",
        ["종목명 및 종목코드", "현재가, 전일대비 등락 금액 및 등락률(%)", "하트 아이콘 (관심 종목 추가/해제 버튼)"],
        ["화면 상단 중앙의 검색창에서 기업명을 검색합니다.", "표출된 가격의 색상(빨강=상승, 파랑=하락)과 수치를 확인합니다.", "계속 지켜보고 싶은 종목이면 하트 아이콘을 눌러 마이페이지에 저장합니다."], "🎯")

    create_slide(prs, "종목 분석 > 기업소개 탭", "종목 분석 > 기업소개 탭",
        "이 회사가 무엇을 만들고 어떻게 돈을 버는지, 기본 정보를 확인합니다.",
        ["기업 개요 요약 (비즈니스 모델 설명)", "시가총액, 52주 최고/최저가 등 핵심 정보 요약", "업종(섹터) 및 속한 테마 정보"],
        ["기업소개 탭을 눌러 이 회사가 구체적으로 어떤 사업을 하는지 짧은 요약을 읽어봅니다."], "📖")

    create_slide(prs, "종목 분석 > 호가 / 실시간 시세 탭", "종목 분석 > 호가 탭 / 실시간 일별시세 탭",
        "현재 시장에서 사람들이 얼마에 사고 팔려고 하는지(호가), 그리고 매일매일 가격이 어떻게 변했는지 확인합니다.",
        ["호가창: 매수(사는 사람)와 매도(파는 사람) 주문이 걸려있는 가격대와 수량 바 그래프", "일별시세: 최근 날짜별 종가, 시가, 거래량 표"],
        ["호가창 탭에서 바 그래프가 긴 쪽을 보며 매수/매도 벽이 어디에 있는지 파악합니다.", "일별시세 탭에서 최근 며칠간 거래량이 터진 날이 언제인지 확인합니다."], "📉")

    create_slide(prs, "종목 분석 > 실적 / 재무분석 / 배당 탭", "종목 분석 > 실적 / 재무 / 배당 탭",
        "회사가 실제로 돈을 잘 벌고 있는지(실적/재무), 그리고 주주들에게 이익을 잘 나눠주는지(배당) 분석합니다.",
        ["실적: 매출액, 영업이익 차트 및 컨센서스(시장 예상치) 달성 여부", "재무: 부채비율, 유보율 등 안정성 지표", "배당: 배당수익률, 과거 배당 지급 이력"],
        ["실적 탭에서 그래프가 매년/매분기 우상향하고 있는지 확인합니다.", "배당 탭에서 배당금을 꾸준히 줬는지(지급 횟수)와 배당률을 확인해 안정적인 투자가 가능한지 판단합니다."], "💰")

    create_slide(prs, "종목 분석 > 뉴스 / 공시 탭", "종목 분석 > 뉴스 / 공시 탭",
        "이 회사와 관련된 최신 기사들과, 회사가 공식적으로 발표한 문서(공시)들을 모아봅니다.",
        ["뉴스: 시간순 뉴스 리스트 및 긍정/부정 감성 아이콘", "공시: 전자공시시스템(DART) 기반 중요 발표 문서 (실적발표, 계약 체결 등)"],
        ["뉴스 탭에서 오늘 뜬 기사가 호재(긍정)인지 악재(부정)인지 배지를 통해 빠르게 훑어봅니다.", "공시 탭에서 '단일판매/공급계약' 등 주가에 영향을 미치는 큰 공시가 없는지 확인합니다."], "📰")

    create_slide(prs, "종목 분석 > 투자자 / 외국인 동향 / 공매도 탭", "종목 분석 > 투자자동향 / 외국인기관 / 공매도 탭",
        "개인, 외국인, 기관 중 누가 이 주식을 주도해서 사고 있는지 수급을 파악합니다.",
        ["투자자 동향: 개인/외인/기관별 일일 순매수 수치표", "공매도: 주가가 떨어질 것에 베팅한 공매도 거래 비중 및 잔고"],
        ["외국인/기관 탭에서 메이저 투자자들이 이 주식을 며칠 연속으로 사모으고 있는지(빨간불) 확인합니다.", "공매도 비중이 갑자기 높아졌다면 단기 하락 압력이 있을 수 있으므로 주의합니다."], "👥")

    # --- 5. 상승탄력 / 관계망 ---
    create_slide(prs, "5. 상승탄력 & 관계망", "분석 툴", "섹션 커버", [], [], "🚀", is_section=True)

    create_slide(prs, "상승탄력 분석", "메뉴 > 상승탄력 분석 (/analysis)",
        "현재 주가가 얼마나 힘 있게 오르고 있는지(모멘텀)를 분석하여 매매 타이밍을 잡도록 도와줍니다.",
        ["상승탄력 점수 (0~100점)", "단기/중기 이동평균선(추세선) 크로스 차트", "과열/침체 지표 바"],
        ["점수가 80점 이상으로 과열 구간에 진입했다면 신규 매수보다는 관망이나 차익실현을 고려합니다.", "단기 선이 중기 선을 위로 뚫고 올라가는 교차점(골든크로스)을 확인하면 매수 타이밍으로 참고합니다."], "📈")

    create_slide(prs, "관계망 (온톨로지)", "메뉴 > 관계망 (/ontology)",
        "주식 시장의 기업과 테마들이 거미줄처럼 얽혀 영향을 주고받는 모습을 시각 지도로 보여줍니다.",
        ["인터랙티브 노드(점) 네트워크 맵", "특정 노드 선택 시 연결된 파급효과 사이드 패널", "확대/축소 및 드래그 캔버스"],
        ["마우스 휠로 맵을 확대하고, 흥미로운 키워드(예: 2차전지 소재) 노드를 클릭해봅니다.", "오른쪽 패널에 뜨는 '이 테마가 오를 때 같이 오를 확률이 높은 기업들' 목록을 확인하여 새로운 투자 아이디어를 얻습니다."], "🕸️")

    # --- 6. 커뮤니티 ---
    create_slide(prs, "6. 커뮤니티", "커뮤니티 /community", "섹션 커버", [], [], "💬", is_section=True)

    create_slide(prs, "커뮤니티 게시판", "메뉴 > 커뮤니티 (/community)",
        "다른 투자자들과 종목에 대한 의견과 정보, 시장 전망을 자유롭게 나누는 공간입니다.",
        ["인기글 토글 탭 (오늘의 인기글 모아보기)", "게시글 목록 (제목, 작성자, 조회수, 좋아요, 댓글 수 표출)", "우측 하단 플로팅 글쓰기 버튼"],
        ["목록에서 흥미로운 제목의 글을 클릭해 내용을 읽고, 하단에 댓글을 달아 의견을 나눕니다.", "우측 상단 돋보기 아이콘으로 특정 종목(예: '카카오')을 검색해 해당 종목 주주들의 생각을 살펴봅니다."], "📋")

    create_slide(prs, "[팝업] 글쓰기 모달", "커뮤니티 > 글쓰기 버튼 클릭",
        "나의 투자 아이디어나 질문을 게시판에 등록하기 위해 텍스트를 작성하는 팝업 화면입니다.",
        ["카테고리 선택 드롭다운 (일반/분석/질문 등)", "관련 종목 태그 추가 입력창", "제목 및 본문 입력 폼", "등록 버튼"],
        ["우측 하단의 동그란 펜 모양 버튼을 누릅니다.", "특정 주식에 대한 이야기라면 '종목 태그'에 이름을 검색해 추가합니다.", "내용을 작성하고 '등록'을 눌러 사람들과 공유합니다."], "📝")

    # --- 7. 마이페이지 / 구독 ---
    create_slide(prs, "7. 마이페이지 및 구독", "개인화 /mypage", "섹션 커버", [], [], "👤", is_section=True)

    create_slide(prs, "마이페이지 & 알림 설정", "메뉴 > 마이페이지 (/mypage)",
        "내 계정 정보를 관리하고, 하트를 누른 관심 종목들과 시스템 알림 수신 여부를 설정합니다.",
        ["내 프로필 정보 (이메일, 닉네임)", "관심 종목 리스트 탭 (삭제 관리 가능)", "알림 설정 탭 (불확실성 알림, 지정가 도달 알림 토글 스위치)"],
        ["관심 종목 탭에서 더 이상 보지 않을 주식은 'X'나 휴지통 아이콘을 눌러 리스트에서 지웁니다.", "알림 탭에서 '시장 급락 경고 알림' 스위치를 켜서 위험 상황 시 팝업 안내를 받도록 설정합니다."], "⚙️")

    create_slide(prs, "구독 (요금제 결제)", "메뉴 > 구독 (/subscription)",
        "무료 버전에서 제한된 기능을 모두 사용하기 위해 프로/프리미엄 요금제로 업그레이드하는 화면입니다.",
        ["플랜별 혜택 비교 표 (무료 / Pro / Premium)", "월간/연간 결제 토글 스위치 (연간 할인율 표시)", "결제하기 버튼 및 현재 이용 중인 플랜 표시"],
        ["월간 결제와 연간 결제 탭을 눌러 가격 차이를 비교해봅니다.", "각 요금제 카드 아래의 체크리스트를 통해 나에게 필요한 기능(예: 실시간 알림 등)이 포함된 요금제를 선택하고 결제합니다."], "💳")

    prs.save('client/public/StockLink_Full_Manual.pptx')

create_manual()
