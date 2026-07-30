from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def create_slide(prs, title, path, purpose, content, usage, icon_text):
    slide_layout = prs.slide_layouts[5] # Blank layout with title
    slide = prs.slides.add_slide(slide_layout)
    
    # Background color (dark theme)
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(15, 23, 42) # slate-900

    # Header Bar
    shapes = slide.shapes
    header_bar = shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1)
    )
    header_bar.fill.solid()
    header_bar.fill.fore_color.rgb = RGBColor(30, 41, 59) # slate-800
    header_bar.line.color.rgb = RGBColor(30, 41, 59)

    # Title
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.6))
    tf = txBox.text_frame
    tf.text = f"{icon_text} {title}"
    tf.paragraphs[0].font.size = Pt(28)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255) # white

    # Content Boxes
    left_margin = Inches(0.5)
    width = Inches(9)
    
    # 1. 화면 경로
    txBox = slide.shapes.add_textbox(left_margin, Inches(1.2), width, Inches(0.5))
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = "📍 화면 경로: "
    p.font.bold = True
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(96, 165, 250) # blue-400
    p.text += path
    p.font.bold = False
    p.font.color.rgb = RGBColor(203, 213, 225) # slate-300

    # 2. 화면 목적
    txBox = slide.shapes.add_textbox(left_margin, Inches(1.8), width, Inches(0.8))
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = "🎯 화면 목적 / 설명"
    p.font.bold = True
    p.font.size = Pt(16)
    p.font.color.rgb = RGBColor(255, 255, 255)
    
    p2 = tf.add_paragraph()
    p2.text = purpose
    p2.font.size = Pt(14)
    p2.font.color.rgb = RGBColor(203, 213, 225)
    p2.space_before = Pt(6)

    # 3. 표출 내용
    txBox = slide.shapes.add_textbox(left_margin, Inches(2.8), Inches(4.3), Inches(4))
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = "📋 화면에 보이는 내용"
    p.font.bold = True
    p.font.size = Pt(16)
    p.font.color.rgb = RGBColor(255, 255, 255)
    
    for item in content:
        p_item = tf.add_paragraph()
        p_item.text = f"• {item}"
        p_item.font.size = Pt(13)
        p_item.font.color.rgb = RGBColor(203, 213, 225)
        p_item.space_before = Pt(4)

    # 4. 사용 방법
    txBox = slide.shapes.add_textbox(Inches(5.2), Inches(2.8), Inches(4.3), Inches(4))
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = "💡 이렇게 사용하세요"
    p.font.bold = True
    p.font.size = Pt(16)
    p.font.color.rgb = RGBColor(255, 255, 255)
    
    for i, item in enumerate(usage):
        p_item = tf.add_paragraph()
        p_item.text = f"{(i+1)}️⃣ {item}"
        p_item.font.size = Pt(13)
        p_item.font.color.rgb = RGBColor(203, 213, 225)
        p_item.space_before = Pt(6)

def create_manual():
    prs = Presentation()
    # 16:9 ratio
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)

    # 1. Title Slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(15, 23, 42)
    
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "StockLink 사용자 매뉴얼"
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    title.text_frame.paragraphs[0].font.size = Pt(44)
    title.text_frame.paragraphs[0].font.bold = True
    
    subtitle.text = "처음 오신 분들을 위한 쉽고 친절한 안내서\n\n시장 데이터부터 AI 분석까지, StockLink 100% 활용하기"
    subtitle.text_frame.paragraphs[0].font.color.rgb = RGBColor(148, 163, 184) # slate-400
    
    # 2. 대시보드 (홈)
    create_slide(prs, 
        "홈 (대시보드)", 
        "홈 (/) > 대시보드",
        "StockLink에 접속하면 가장 먼저 만나는 화면으로, 현재 주식 시장의 흐름과 내가 관심 있는 종목들의 상태를 한눈에 파악할 수 있습니다.",
        [
            "주요 지수 요약: KOSPI, 코스닥 등 주요 지표의 현재 수치",
            "글로벌 증시 요약: 해외 주요 시장의 전반적인 분위기",
            "상승/하락 종목: 오늘 가장 눈에 띄게 오르거나 내린 주식 목록",
            "내 관심 종목: 내가 미리 추가해둔 주식들의 최신 현황"
        ],
        [
            "오늘 시장의 전체적인 분위기(빨간색이 많으면 상승장, 파란색이 많으면 하락장)를 먼저 확인합니다.",
            "상단 요약 카드를 통해 오늘 특별히 많이 오르거나 내린 주식들을 살펴봅니다.",
            "더 자세히 알고 싶은 지표나 종목이 있다면 클릭해서 상세 페이지로 이동합니다."
        ],
        "🏠"
    )

    # 3. 시장 현황
    create_slide(prs, 
        "시장 현황 (Market Analysis)", 
        "메뉴 > 시장 현황 (/market)",
        "국내외 시장 전체가 현재 어떤 방향으로 흘러가고 있는지 깊이 있게 분석해주는 화면입니다.",
        [
            "시장 신호등: 시장의 과열(너무 비쌈) 또는 침체(너무 쌈) 상태를 알려주는 지표",
            "투자자 동향: 기관, 외국인, 개인 투자자들이 지금 주식을 사고 있는지 팔고 있는지 보여주는 차트",
            "주요 뉴스 요약: 시장에 영향을 미치는 오늘의 중요 뉴스"
        ],
        [
            "상단의 '시장 신호등'을 보고 지금 투자를 공격적으로 할지 보수적으로 할지 감을 잡습니다.",
            "외국인과 기관이 어떤 방향으로 움직이는지 그래프를 통해 확인합니다.",
            "오늘 시장을 움직인 주요 뉴스의 헤드라인들을 훑어봅니다."
        ],
        "📊"
    )

    # 4. 종목 분석
    create_slide(prs, 
        "종목 상세 분석 (Stock)", 
        "메뉴 > 종목 분석 (/stock)",
        "특정 주식(예: 삼성전자)을 검색하면, 해당 주식에 대한 과거 기록, 현재 가격, 그리고 앞으로의 전망까지 모든 정보를 확인할 수 있습니다.",
        [
            "기본 정보: 현재가, 등락폭(오르내린 금액), 거래량",
            "주가 차트: 과거부터 지금까지의 주가 흐름을 보여주는 선 그래프",
            "기업 가치: 이 회사가 돈을 얼마나 잘 버는지(실적) 요약한 표",
            "최근 이슈: 이 주식과 관련된 최근 뉴스나 이벤트"
        ],
        [
            "상단 검색창에 궁금한 기업의 이름이나 번호(종목코드)를 입력합니다.",
            "가장 위에 표시된 현재가와 오르내린 폭을 확인합니다.",
            "아래로 스크롤하여 이 기업의 최신 뉴스나 실적 요약을 읽어봅니다."
        ],
        "🏢"
    )

    # 5. 소셜 분석
    create_slide(prs, 
        "소셜 분석 (Social Analysis)", 
        "메뉴 > 인사이트 > 소셜 분석 (/social 또는 /insight)",
        "유명한 투자자나 경제 전문가, 혹은 SNS(트위터 등)에서 주식에 대해 어떤 이야기들이 오가고 있는지 모아서 보여주는 화면입니다.",
        [
            "주요 하이라이트: 오늘 전문가들이 가장 중요하게 이야기한 내용 요약",
            "긍정 영향(📈): 전문가 발언으로 인해 상승탄력을 받을 것으로 기대되는 종목들",
            "부정 영향(📉): 악재로 인해 하락방어가 필요한, 우려되는 종목들",
            "섹터별 분위기: 반도체, 2차전지 등 분야별로 긍정적인 이야기가 많은지 파이 차트로 표출"
        ],
        [
            "우측 상단의 날짜 선택기를 눌러 원하는 날짜의 리포트를 불러옵니다.",
            "아래쪽 종합 요약 테이블에서 유명 인사들의 발언 요약을 훑어봅니다.",
            "관심 있는 인물의 행을 클릭하면 나타나는 팝업(모달)에서 긍정/부정 종목과 발언 원문을 자세히 읽어봅니다."
        ],
        "🗣️"
    )

    # 6. 온톨로지
    create_slide(prs, 
        "관계망 (Ontology)", 
        "메뉴 > 관계망 (/ontology)",
        "주식 시장의 다양한 기업, 산업, 뉴스들이 서로 어떻게 연결되어 영향을 주고받는지 거미줄 같은 지도로 보여줍니다.",
        [
            "관계망 지도: 점(노드)과 선(링크)으로 이루어진 복잡한 지도",
            "테마 그룹: 같은 이슈로 묶이는 주식들 (예: AI 관련주 묶음)",
            "파급 효과 예측: 한 기업의 문제가 다른 기업에 미칠 영향"
        ],
        [
            "화면 중앙의 거미줄 같은 지도에서 가장 크고 눈에 띄는 원(키워드)을 찾습니다.",
            "해당 원을 클릭하면 그와 연결된 다른 주식들이나 뉴스들이 하이라이트됩니다.",
            "오른쪽 패널에 나오는 상세 설명을 읽어보며, 이 기업이 왜 저 기업과 연결되는지 이해합니다."
        ],
        "🕸️"
    )

    # 7. 상승탄력 분석
    create_slide(prs, 
        "상승탄력 분석 (Analysis / Momentum)", 
        "메뉴 > 상승탄력 분석 (/analysis)",
        "현재 주가가 얼마나 힘 있게 오르고 있는지, 혹은 내리고 있는지 그 추세(힘)를 분석하여 투자 타이밍을 잡도록 도와줍니다.",
        [
            "탄력 점수: 현재 주가가 오르려는 힘이 강한지 점수로 표현 (0~100점)",
            "단기/중기 추세선: 짧은 기간과 중간 기간 동안의 주가 방향성",
            "과열/침체 지표: 지금 사기엔 너무 비싼지(과열), 사기 좋은 때인지(침체) 알려주는 신호"
        ],
        [
            "가장 큰 숫자인 '탄력 점수'를 확인합니다. 점수가 높을수록 상승하는 힘이 강합니다.",
            "단기 추세선이 중기 추세선을 위로 뚫고 올라가는지(좋은 신호) 확인합니다.",
            "너무 과열된 상태라면 잠시 매수를 보류하고 하락방어에 집중합니다."
        ],
        "🚀"
    )

    # 8. 커뮤니티
    create_slide(prs, 
        "커뮤니티 (Community)", 
        "메뉴 > 커뮤니티 (/community)",
        "StockLink를 사용하는 다른 사람들과 투자 아이디어를 나누고 토론할 수 있는 공간입니다.",
        [
            "실시간 인기글: 지금 사람들이 가장 많이 읽고 댓글을 단 게시물",
            "종목 토론방: 특정 주식에 대해 이야기 나누는 공간",
            "글쓰기 버튼: 내 생각이나 질문을 올릴 수 있는 기능"
        ],
        [
            "관심 있는 제목의 글을 클릭해 다른 사람들의 의견을 읽어봅니다.",
            "동의하거나 다른 생각이 있다면 하단에 댓글을 남깁니다.",
            "오른쪽 상단의 '글쓰기' 버튼을 눌러 나의 투자 아이디어를 공유합니다."
        ],
        "💬"
    )

    # 9. 마이페이지 & 구독
    create_slide(prs, 
        "마이페이지 & 구독 (My Page / Subscription)", 
        "메뉴 > 마이페이지 (/mypage) 또는 구독 (/subscription)",
        "내 정보, 설정, 관심 종목 리스트를 관리하고, 더 많은 기능을 쓰기 위해 유료 구독 요금제를 확인할 수 있습니다.",
        [
            "내 프로필: 가입된 이메일과 닉네임",
            "관심 종목 설정: 내가 하트를 누른 주식 목록 관리",
            "알림 설정: 중요한 뉴스가 뜰 때 알람을 받을지 선택",
            "요금제 안내: 무료/프로/프리미엄 요금제의 혜택 비교 표"
        ],
        [
            "관심 종목 탭에서 더 이상 보지 않을 주식은 '삭제' 버튼을 눌러 정리합니다.",
            "나에게 필요한 알림(예: 급락 경고 등)만 오도록 스위치를 껐다 켭니다.",
            "구독 페이지에서 나에게 맞는 요금제를 비교해보고 필요하다면 업그레이드를 진행합니다."
        ],
        "👤"
    )

    prs.save('StockLink_User_Manual.pptx')

create_manual()
